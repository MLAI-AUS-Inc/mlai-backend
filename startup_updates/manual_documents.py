from __future__ import annotations

import csv
import io
import zipfile
from dataclasses import dataclass
from html import unescape
from pathlib import PurePosixPath
from xml.etree import ElementTree


@dataclass
class ParsedManualDocument:
    extracted_text: str
    extraction_status: str
    parse_notes: str
    last_error: str = ""


def _missing_dependency(name: str) -> ParsedManualDocument:
    return ParsedManualDocument("", "error", "parser_dependency_missing", name)


def _text_from_pdf(raw_bytes: bytes) -> ParsedManualDocument:
    try:
        from pypdf import PdfReader
    except ImportError:
        return _missing_dependency("pypdf")

    reader = PdfReader(io.BytesIO(raw_bytes))
    pages = [page.extract_text() or "" for page in reader.pages]
    return ParsedManualDocument("\n".join(item.strip() for item in pages if item.strip()), "processed", "pdf_parsed")


def _text_from_docx(raw_bytes: bytes) -> ParsedManualDocument:
    try:
        from docx import Document
    except ImportError:
        return _text_from_zipped_xml(raw_bytes, "word/document.xml", "docx_xml_parsed")

    doc = Document(io.BytesIO(raw_bytes))
    text = "\n".join(paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text.strip())
    return ParsedManualDocument(text, "processed", "docx_parsed")


def _text_from_pptx(raw_bytes: bytes) -> ParsedManualDocument:
    try:
        from pptx import Presentation
    except ImportError:
        return _text_from_zipped_xml(raw_bytes, "ppt/slides/slide", "pptx_xml_parsed", prefix=True)

    presentation = Presentation(io.BytesIO(raw_bytes))
    lines = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                lines.append(text)
    return ParsedManualDocument("\n".join(lines), "processed", "pptx_parsed")


def _text_from_xlsx(raw_bytes: bytes) -> ParsedManualDocument:
    try:
        import openpyxl
    except ImportError:
        return _missing_dependency("openpyxl")

    workbook = openpyxl.load_workbook(io.BytesIO(raw_bytes), read_only=True, data_only=True)
    lines = []
    for sheet in workbook.worksheets:
        lines.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if cells:
                lines.append(" | ".join(cells))
    return ParsedManualDocument("\n".join(lines), "processed", "xlsx_parsed")


def _text_from_csv(raw_bytes: bytes) -> ParsedManualDocument:
    stream = io.StringIO(raw_bytes.decode("utf-8", errors="ignore"))
    reader = csv.reader(stream)
    text = "\n".join(" | ".join(cell.strip() for cell in row if cell.strip()) for row in reader if row)
    return ParsedManualDocument(text, "processed", "csv_parsed")


def _text_from_html(raw_bytes: bytes) -> ParsedManualDocument:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        text = raw_bytes.decode("utf-8", errors="ignore")
        return ParsedManualDocument(unescape(text), "processed", "html_raw_text_parsed")

    soup = BeautifulSoup(raw_bytes.decode("utf-8", errors="ignore"), "html.parser")
    return ParsedManualDocument(unescape(soup.get_text("\n", strip=True)), "processed", "html_parsed")


def _text_from_rtf(raw_bytes: bytes) -> ParsedManualDocument:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        return _missing_dependency("striprtf")

    return ParsedManualDocument(
        rtf_to_text(raw_bytes.decode("utf-8", errors="ignore")).strip(),
        "processed",
        "rtf_parsed",
    )


def _text_from_odt(raw_bytes: bytes) -> ParsedManualDocument:
    return _text_from_zipped_xml(raw_bytes, "content.xml", "odt_parsed")


def _text_from_zipped_xml(raw_bytes: bytes, member: str, parse_notes: str, *, prefix: bool = False) -> ParsedManualDocument:
    lines = []
    with zipfile.ZipFile(io.BytesIO(raw_bytes)) as archive:
        names = archive.namelist()
        selected_names = [name for name in names if name.startswith(member)] if prefix else [member]
        for name in selected_names:
            if name not in names or not name.endswith(".xml"):
                continue
            root = ElementTree.fromstring(archive.read(name))
            lines.extend(text.strip() for text in root.itertext() if text and text.strip())
    return ParsedManualDocument("\n".join(lines), "processed", parse_notes)


def parse_manual_document(*, filename: str, content_type: str, raw_bytes: bytes) -> ParsedManualDocument:
    if not raw_bytes:
        return ParsedManualDocument("", "error", "empty_document", "")

    extension = PurePosixPath(filename or "").suffix.lower()
    mime_type = str(content_type or "").split(";")[0].strip().lower()

    try:
        if extension == ".pdf" or mime_type == "application/pdf":
            return _text_from_pdf(raw_bytes)
        if extension == ".docx" or mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return _text_from_docx(raw_bytes)
        if extension == ".pptx" or mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return _text_from_pptx(raw_bytes)
        if extension in {".xlsx", ".xlsm"} or mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return _text_from_xlsx(raw_bytes)
        if extension == ".csv" or mime_type == "text/csv":
            return _text_from_csv(raw_bytes)
        if extension in {".html", ".htm"} or mime_type == "text/html":
            return _text_from_html(raw_bytes)
        if extension == ".rtf" or mime_type in {"application/rtf", "text/rtf"}:
            return _text_from_rtf(raw_bytes)
        if extension == ".odt" or mime_type == "application/vnd.oasis.opendocument.text":
            return _text_from_odt(raw_bytes)
        if extension in {".md", ".txt"} or mime_type.startswith("text/"):
            return ParsedManualDocument(
                raw_bytes.decode("utf-8", errors="ignore").strip(),
                "processed",
                "text_parsed",
            )
        if mime_type.startswith("image/"):
            return ParsedManualDocument("", "unsupported", "ocr_not_supported", "")
        return ParsedManualDocument("", "unsupported", "unsupported_document_type", "")
    except Exception as exc:
        return ParsedManualDocument("", "error", f"parse_error:{type(exc).__name__}", str(exc)[:500])
